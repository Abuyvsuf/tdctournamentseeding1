"""
build_presentation_v2.py
-------------------------
Slots coded/pooled team lists into a region's PPTX template.

Unlike the first version, this does NOT assume slide order or a fixed
pool count. For every slide, it looks at ALL the text on that slide to
work out which (language, category) it belongs to:

    "Debate: Grade 10 ..."   -> ("English", "Junior")
    "Debate: Senior ..."     -> ("English", "Senior")
    "Mjadala: Grade 10 ..."  -> ("Kiswahili", "Junior")
    "Mjadala: Senior ..."    -> ("Kiswahili", "Senior")

Then it finds every "Pool X Teams ... so far ..." shape on that slide
(however many there are -- 2, 4, whatever the template actually has) and
fills each with the matching pool's teams. Any slide that doesn't match
one of these patterns (cover, about, public speaking, qualifiers, thank
you, etc.) is left completely untouched.
"""

import copy
import re

from pptx import Presentation

POOL_HEADER_RE = re.compile(r"Pool\s+([A-Z])\s+Teams.*so far", re.IGNORECASE)

SLIDE_CATEGORY_PATTERNS = [
    (re.compile(r"debate:\s*grade\s*10", re.IGNORECASE), "English", "Junior"),
    (re.compile(r"debate:\s*senior", re.IGNORECASE), "English", "Senior"),
    (re.compile(r"mjadala:\s*grade\s*10", re.IGNORECASE), "Kiswahili", "Junior"),
    (re.compile(r"mjadala:\s*senior", re.IGNORECASE), "Kiswahili", "Senior"),
]


def _slide_all_text(slide):
    parts = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            for p in shape.text_frame.paragraphs:
                if p.text.strip():
                    parts.append(p.text.strip())
    return "\n".join(parts)


def _identify_slide_category(slide_text):
    for pattern, language, category in SLIDE_CATEGORY_PATTERNS:
        if pattern.search(slide_text):
            return language, category
    return None


def _set_paragraph_text(p_element, text, a_ns, font_size_pt=None):
    runs = p_element.findall(f"{{{a_ns}}}r")
    if not runs:
        return
    first_run = runs[0]
    t_el = first_run.find(f"{{{a_ns}}}t")
    if t_el is None:
        t_el = first_run.makeelement(f"{{{a_ns}}}t", {})
        first_run.append(t_el)
    t_el.text = text
    for extra in runs[1:]:
        p_element.remove(extra)
    if font_size_pt is not None:
        rpr = first_run.find(f"{{{a_ns}}}rPr")
        if rpr is not None:
            rpr.set("sz", str(int(font_size_pt * 100)))


# The template's pool boxes were laid out to comfortably fit this many
# team lines at full size. Beyond that, shrink font size proportionally
# (same idea as PowerPoint's own autofit) rather than letting text overflow
# the slide. Floor prevents it from becoming illegible.
BASELINE_TEAM_COUNT = 8
BASE_FONT_PT = 12
FLOOR_FONT_PT = 7


def _font_size_for_count(count):
    if count <= BASELINE_TEAM_COUNT:
        return BASE_FONT_PT
    scale = BASELINE_TEAM_COUNT / count
    size = max(FLOOR_FONT_PT, round(BASE_FONT_PT * scale))
    return size


def _rebuild_multi_paragraph_style(tf, header_p, content_paragraphs, team_names, font_size, a_ns):
    """Style A: each team is already its own <a:p> (e.g. Metropolitan deck).
    Clone the first content paragraph's XML as a formatting template, delete
    the old ones, insert one new <a:p> per team."""
    template_p_xml = copy.deepcopy(content_paragraphs[0]._p)
    body = tf._txBody
    for p in content_paragraphs:
        body.remove(p._p)

    insert_after = header_p._p
    for team_name in team_names:
        new_p = copy.deepcopy(template_p_xml)
        _set_paragraph_text(new_p, team_name, a_ns, font_size_pt=font_size)
        insert_after.addnext(new_p)
        insert_after = new_p


def _rebuild_single_paragraph_style(blob_paragraph, team_names, font_size, a_ns):
    """Style B: the whole team list lives inside ONE <a:p>, with manual
    "1.", "2." numbering typed into the text and <a:br> elements separating
    lines (e.g. Upper Eastern deck). Rebuild that single paragraph's run
    list from scratch, reusing its first run's formatting as a template."""
    p_el = blob_paragraph._p
    runs = p_el.findall(f"{{{a_ns}}}r")
    if not runs:
        return
    template_rpr = runs[0].find(f"{{{a_ns}}}rPr")

    # Remove every existing run and line-break, keep pPr.
    for child in list(p_el):
        tag = child.tag.split("}")[-1]
        if tag in ("r", "br"):
            p_el.remove(child)

    for i, team_name in enumerate(team_names):
        r_el = p_el.makeelement(f"{{{a_ns}}}r", {})
        if template_rpr is not None:
            rpr_copy = copy.deepcopy(template_rpr)
            if font_size is not None:
                rpr_copy.set("sz", str(int(font_size * 100)))
            if "err" in rpr_copy.attrib:
                del rpr_copy.attrib["err"]
            r_el.append(rpr_copy)
        t_el = r_el.makeelement(f"{{{a_ns}}}t", {})
        t_el.text = f"{i + 1}.{team_name}"
        r_el.append(t_el)
        p_el.append(r_el)

        if i < len(team_names) - 1:
            br_el = p_el.makeelement(f"{{{a_ns}}}br", {})
            if template_rpr is not None:
                br_rpr = copy.deepcopy(template_rpr)
                if font_size is not None:
                    br_rpr.set("sz", str(int(font_size * 100)))
                if "err" in br_rpr.attrib:
                    del br_rpr.attrib["err"]
                br_el.append(br_rpr)
            p_el.append(br_el)


def build_presentation(template_path, output_path, pools_by_language_category):
    """
    pools_by_language_category: {
        ("English", "Junior"): {"A": [...], "B": [...], "C": [...], "D": [...]},
        ("English", "Senior"): {...},
        ("Kiswahili", "Junior"): {"A": [...], "B": [...]},
        ("Kiswahili", "Senior"): {...},
    }
    Returns a list of warning strings (e.g. a slide's pool count didn't
    match the data's pool count) for the caller to surface.
    """
    prs = Presentation(template_path)
    a_ns = "http://schemas.openxmlformats.org/drawingml/2006/main"
    warnings = []
    matched_keys = set()

    for slide in prs.slides:
        slide_text = _slide_all_text(slide)
        key = _identify_slide_category(slide_text)
        if key is None:
            continue
        if key not in pools_by_language_category:
            continue
        matched_keys.add(key)
        pools = pools_by_language_category[key]

        pool_shapes = []
        for shape in slide.shapes:
            if not shape.has_text_frame or not shape.text_frame.paragraphs:
                continue
            header_text = shape.text_frame.paragraphs[0].text
            m = POOL_HEADER_RE.search(header_text)
            if m:
                pool_shapes.append((shape, m.group(1)))

        if len(pool_shapes) != len(pools):
            warnings.append(
                f"{key[0]} {key[1]}: template slide has {len(pool_shapes)} pool box(es) "
                f"but data has {len(pools)} pool(s) -- using whichever pool letters match."
            )

        for shape, pool_letter in pool_shapes:
            if pool_letter not in pools:
                continue
            tf = shape.text_frame
            paragraphs = tf.paragraphs
            header_p = paragraphs[0]

            trailing = []
            for p in reversed(paragraphs):
                if len(p.runs) == 0 and len(p._p.findall(f"{{{a_ns}}}br")) == 0 and p is not header_p:
                    trailing.append(p)
                else:
                    break
            trailing_count = len(trailing)

            content_paragraphs = (
                paragraphs[1: len(paragraphs) - trailing_count] if trailing_count else paragraphs[1:]
            )
            if not content_paragraphs:
                continue

            team_names = pools[pool_letter]
            font_size = _font_size_for_count(len(team_names))

            # Detect structure: does any content paragraph contain <a:br>
            # line breaks (the "whole list packed into one paragraph" style)?
            blob_paragraphs = [p for p in content_paragraphs if len(p._p.findall(f"{{{a_ns}}}br")) > 0]

            if blob_paragraphs:
                _rebuild_single_paragraph_style(blob_paragraphs[0], team_names, font_size, a_ns)
                # if there happened to be more than one blob paragraph (rare),
                # empty out any extras so old data doesn't linger
                for extra in blob_paragraphs[1:]:
                    _rebuild_single_paragraph_style(extra, [], font_size, a_ns)
            else:
                _rebuild_multi_paragraph_style(tf, header_p, content_paragraphs, team_names, font_size, a_ns)

    unmatched = set(pools_by_language_category.keys()) - matched_keys
    for language, category in sorted(unmatched):
        team_count = sum(len(t) for t in pools_by_language_category[(language, category)].values())
        if team_count == 0:
            continue  # no real teams in this bucket anyway, nothing lost
        warnings.append(
            f"{language} {category}: this template has NO matching slide -- "
            f"{team_count} team(s) of coded/pooled data could not be placed "
            f"anywhere in this presentation. Add a slide titled like "
            f"'{'Debate' if language == 'English' else 'Mjadala'}: "
            f"{'Grade 10' if category == 'Junior' else 'Senior'}' with Pool boxes, "
            f"or this region truly has no {language} {category} participants."
        )

    prs.save(output_path)
    return warnings
